import sys
from datetime import datetime
from pathlib import Path
from xml.sax.saxutils import escape

from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

from pptx import Presentation

# Generated pdf or ppt outputs desktop
OUTPUT_DIR = Path.home() / "Desktop"


# Makes unique pdf and ppt names
def _unique_path(ext: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True) 
    ts = datetime.now().strftime("%Y-%m-%d_%H%M%S")
    return OUTPUT_DIR / f"report_{ts}.{ext}"


# Extract heading and body from a section, with some flexibility in field names and formats.
def _section_parts(section) -> tuple[str, str]:
    if isinstance(section, str):
        return "", section
    if not isinstance(section, dict):
        return "", str(section)
    heading = (section.get("heading") or section.get("title")
               or section.get("name") or section.get("header") or "")
    body = (section.get("body") or section.get("content")
            or section.get("text") or section.get("description") or "")
    if isinstance(body, list):
        body = "\n".join(str(b) for b in body)
    return str(heading), str(body)


def make_pdf(title: str, sections: list[dict]) -> str:
    path = _unique_path("pdf") # output file path with extension pdf
    doc = SimpleDocTemplate(str(path), pagesize=LETTER) 
    styles = getSampleStyleSheet()

    flowables = [
        Paragraph(escape(title), styles["Title"]),
        Spacer(1, 0.3 * inch),
    ]

    # append the sections into flowables
    for section in sections:
        heading, body = _section_parts(section)
        flowables.append(Paragraph(escape(heading), styles["Heading2"]))
        flowables.append(Paragraph(escape(body), styles["BodyText"]))
        flowables.append(Spacer(1, 0.2 * inch))

    doc.build(flowables) # write the flowables into the document
    return str(path)  


def make_pptx(title: str, sections: list[dict]) -> str:
    path = _unique_path("pptx") # output file with extension pptx
    prs = Presentation() # create a new presentation

    # add a title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title

    # add a slide for each section, with the heading as the title and body as bullet points
    for section in sections:
        heading, body = _section_parts(section)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = heading

        lines = [ln.strip() for ln in body.split("\n") if ln.strip()]

        tf = slide.placeholders[1].text_frame
        tf.text = lines[0] if lines else ""
        for line in lines[1:]:
            p = tf.add_paragraph()
            p.text = line

    prs.save(str(path))
    return str(path)
